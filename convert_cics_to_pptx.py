#!/usr/bin/env python3
"""
Convert CICS/CICS-new.docx to CICS_Presentation.pptx

Rules:
- 100% content preserved (no skipping)
- Max 6-8 bullets per slide OR ~450 chars of body text
- Code blocks in separate slides with monospace font
- Split slides named "Title (Contd. N)"
- No text overflow
"""

import re
import os
import textwrap
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Layout constants ──────────────────────────────────────────────────────────
SLIDE_W = Inches(13.33)   # 16:9 widescreen
SLIDE_H = Inches(7.5)

TITLE_TOP    = Inches(0.35)
TITLE_LEFT   = Inches(0.5)
TITLE_WIDTH  = Inches(12.33)
TITLE_HEIGHT = Inches(1.1)

BODY_TOP    = Inches(1.5)
BODY_LEFT   = Inches(0.5)
BODY_WIDTH  = Inches(12.33)
BODY_HEIGHT = Inches(5.7)

# Colours
C_TITLE_BG   = RGBColor(0x1F, 0x39, 0x64)   # dark navy
C_BODY_BG    = RGBColor(0xF2, 0xF2, 0xF2)   # light grey
C_TITLE_TEXT = RGBColor(0xFF, 0xFF, 0xFF)   # white
C_BODY_TEXT  = RGBColor(0x1F, 0x1F, 0x1F)   # near-black
C_CODE_BG    = RGBColor(0x1E, 0x1E, 0x1E)   # dark background for code
C_CODE_TEXT  = RGBColor(0xD4, 0xD4, 0xD4)   # light text for code
C_ACCENT     = RGBColor(0x2E, 0x75, 0xB6)   # accent blue
C_DIVIDER_BG = RGBColor(0x2E, 0x75, 0xB6)   # section divider background

# Content limits
MAX_BULLETS   = 7    # max bullet points per slide
MAX_CHARS     = 450  # max total body chars per slide
MAX_CODE_LINES = 18  # max code lines per slide

# ── Code-block detection ──────────────────────────────────────────────────────
CODE_PATTERNS = [
    re.compile(r'^\s*EXEC\s+CICS', re.IGNORECASE),
    re.compile(r'^\s*END-EXEC', re.IGNORECASE),
    re.compile(r'^\s*(SERVICE\s+RELOAD|PROCEDURE\s+DIVISION|LINKAGE\s+SECTION|WORKING-STORAGE)', re.IGNORECASE),
    re.compile(r'^\s*\d{2}\s+\w+\s+PIC\s+', re.IGNORECASE),
    re.compile(r'^\s*(01|02|03|05|10|15)\s+\w'),
    re.compile(r'^\s*(WHEN|EVALUATE|IF|ELSE|END-IF|END-EVALUATE|PERFORM|GO\s+TO)\s+', re.IGNORECASE),
    re.compile(r'DFHMSD|DFHMDI|DFHMDF|DFHBMSCA', re.IGNORECASE),
    re.compile(r'^\s*(MOVE|ADD|SUBTRACT|MULTIPLY|DIVIDE|COMPUTE)\s+.+\s+TO\s+', re.IGNORECASE),
    re.compile(r'TTTTrr|[A-Z]{8}\s*[-]{1,3}\s*8 bytes'),
    # CICS option/parameter syntax fragments: [OPTION(...)], (arg-value), etc.
    re.compile(r'^\s*\[(?!Note|note|\()'),              # lines starting with [ (but not [Note] etc.)
    re.compile(r'^\s*\((?=[A-Z0-9-]{2,}\)|\d)'),        # lines starting with ( followed by identifier or digit
    re.compile(r'\(data-value\)|\(data-area\)|\(name\)|\(hhmmss\)|\(integer\)', re.IGNORECASE),
    re.compile(r'^\s*(FROM|INTO|SET|LENGTH|FLENGTH|INTERVAL|TIME|TRANSID|TERMID|QUEUE|REQID|SYSID|RIDFLD|KEYLENGTH|RESP|TOKEN|NODE|USERID)\s*\(', re.IGNORECASE),
    re.compile(r'^\s*(CURSOR|ERASE|FREEKB|ALARM|FRSET|ERASEUP|MAPONLY|DATAONLY|ACCUM|PAGING|NLEOM|PRINT|RBA|RRN|EQUAL|GTEQ|GENERIC|UPDATE|TOKEN)\s*$', re.IGNORECASE),
    re.compile(r'^\s*(HEADER|TRAILER|TITLE|LIST|OPCLASS|ERRTERM)\s*\(', re.IGNORECASE),
    re.compile(r'^\s*(RTRANSID|RTERMID|SPOOLOPEN|SPOOLWRITE|SPOOLCLOSE)\s*(\(|OUTPUT|FROM|TOKEN|INPUT|RESP)', re.IGNORECASE),
    re.compile(r'^\s*(ADDRESS\s+OF|EXEC\s+CICS|END-EXEC|SERVICE\s+RELOAD)', re.IGNORECASE),
]

COBOL_KEYWORDS = {
    'PIC', 'COMP', 'REDEFINES', 'VALUE', 'FILLER', 'OCCURS',
    'INDEXED', 'BY', 'DEPENDING', 'SECTION', 'DIVISION',
}

# Words that start a sentence / body text (not a heading)
SENTENCE_STARTERS = re.compile(
    r'^(It |The |This |If |In |When |As |So |But |For |A |An |We |You |To |By |'
    r'On |At |With |From |There |All |Since |Once |After |Before |During |While |'
    r'Although |However |Because |Note |Please |Any |Some |Each |Every |Most |Many |'
    r'Both |Either |Neither |These |Those |Such |Same |Other |Another |Its |Their |'
    r'Our |Your |His |Her )',
    re.IGNORECASE,
)

# Sentence verbs that indicate body text (not a heading title)
SENTENCE_VERBS = re.compile(
    r'\b(is |are |was |were |can |will |should |must |have |has |had |does |do |'
    r'may |might |shall |could |would |allows |permits |provides |specifies |'
    r'contains |defines |returns |sends |receives |opens |closes |reads |writes |'
    r'establishes |determines |names |controls |places |stores |indicates |'
    r'identifies |represents)\b',
    re.IGNORECASE,
)


def is_code_line(text: str) -> bool:
    """Return True if this line looks like COBOL / BMS / CICS code."""
    t = text.strip()
    if not t:
        return False
    for pat in CODE_PATTERNS:
        if pat.search(t):
            return True
    # Lines with multiple ALL-CAPS COBOL keywords
    words = set(re.findall(r'\b[A-Z][A-Z0-9-]{2,}\b', t))
    if len(words & COBOL_KEYWORDS) >= 2:
        return True
    return False


def looks_like_heading(text: str) -> bool:
    """
    Heuristic: return True if a short line looks like a section/sub-section
    heading rather than body text.
    """
    # Must be short
    if len(text) > 72:
        return False

    # Must not end with sentence-ending punctuation
    if text[-1:] in ('.', ',', ';'):
        return False

    # Must not start with a lowercase letter (continuation of prior sentence)
    if text[0].islower():
        return False

    # Must not start with special characters (code/citation/list markers)
    if text[0] in '([{(*#=+~-"\'':
        return False

    # Must not start like a normal English sentence
    if SENTENCE_STARTERS.match(text):
        return False

    # Must not read like body text with a main verb
    if SENTENCE_VERBS.search(text):
        return False

    # Must not look like a numbered list item
    if re.match(r'^\d+[\.\)]\s', text):
        return False

    # Must not be trivially short noise
    if len(text) <= 3:
        return False

    return True


# ── Paragraph classification ──────────────────────────────────────────────────
def classify_paragraphs(doc: Document):
    """
    Return list of dicts:
      {'type': 'TITLE'|'H1'|'H2'|'BODY'|'CODE'|'BLANK', 'text': str}
    """
    rows = []
    for p in doc.paragraphs:
        raw  = p.text
        text = raw.strip()
        style = p.style.name

        if not text or text in (' ', '\xa0'):
            rows.append({'type': 'BLANK', 'text': ''})
            continue

        # Explicit heading styles
        if style == 'Heading 1':
            rows.append({'type': 'H1', 'text': text})
            continue
        if style in ('Heading 2', 'Heading 3'):
            rows.append({'type': 'H2', 'text': text})
            continue
        if style in ('Heading 4', 'Heading 5', 'Heading 6'):
            rows.append({'type': 'H2', 'text': text})
            continue

        # Code detection (before heading heuristics)
        if is_code_line(text):
            rows.append({'type': 'CODE', 'text': text})
            continue

        # Heuristic sub-heading detection
        if looks_like_heading(text):
            rows.append({'type': 'H2', 'text': text})
            continue

        rows.append({'type': 'BODY', 'text': text})

    return rows


def extract_tables(doc: Document):
    """Return all table rows as list of list-of-strings."""
    tables = []
    for tbl in doc.tables:
        rows = []
        for row in tbl.rows:
            cells = [c.text.replace('\n', ' ').strip() for c in row.cells]
            rows.append(cells)
        tables.append(rows)
    return tables


# ── Content grouping ──────────────────────────────────────────────────────────
class Section:
    """A logical section with a title and a list of content items."""
    def __init__(self, title: str, level: int = 1):
        self.title   = title
        self.level   = level   # 1 = major, 2 = sub
        self.items   = []      # list of {'type': 'BODY'|'CODE', 'text': str}

    def add(self, item_type: str, text: str):
        self.items.append({'type': item_type, 'text': text})


def build_sections(rows):
    """
    Convert the flat paragraph list into a hierarchy of Sections.
    """
    sections = []
    current  = Section('CICS – Introduction', level=1)

    # Track previous non-blank type for context
    for i, row in enumerate(rows):
        t    = row['type']
        text = row['text']

        if t == 'BLANK':
            continue

        if t == 'H1':
            if current.items or current.title != 'CICS – Introduction':
                sections.append(current)
            current = Section(text, level=1)

        elif t == 'H2':
            # If this is the very first item and no body yet, use as sub-title
            if not current.items:
                # Sub-section within the current H1
                sections.append(current)
                current = Section(text, level=2)
            else:
                sections.append(current)
                current = Section(text, level=2)

        elif t in ('BODY', 'CODE'):
            current.add(t, text)

    if current.items or sections == []:
        sections.append(current)

    return sections


# ── PPTX construction ─────────────────────────────────────────────────────────
def set_slide_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size,
                 bold=False, color=None, align=PP_ALIGN.LEFT,
                 wrap=True, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.name  = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def add_title_box(slide, title_text):
    """Add the dark-navy title bar at the top of a content slide."""
    # Background rectangle for title
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_TITLE_BG
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    run.font.size  = Pt(22)
    run.font.bold  = True
    run.font.color.rgb = C_TITLE_TEXT
    run.font.name  = 'Calibri'


def split_bullets(items, max_b=MAX_BULLETS, max_c=MAX_CHARS):
    """
    Split a list of body items into groups that fit within one slide.
    Returns list of groups.
    Long individual items are split at sentence boundaries.
    """
    # First, expand any items that are too long on their own
    expanded = []
    for item in items:
        if len(item['text']) > max_c:
            # Split long item at sentence boundaries or word-wrap
            sentences = re.split(r'(?<=[.!?])\s+', item['text'])
            chunk_text = ''
            for s in sentences:
                if chunk_text and len(chunk_text) + len(s) > max_c:
                    expanded.append({'type': item['type'], 'text': chunk_text.strip()})
                    chunk_text = s
                else:
                    chunk_text = (chunk_text + ' ' + s).strip() if chunk_text else s
            if chunk_text:
                expanded.append({'type': item['type'], 'text': chunk_text.strip()})
        else:
            expanded.append(item)

    # Now group into slides
    groups   = []
    current  = []
    cur_chars = 0

    for item in expanded:
        blen = len(item['text'])
        if (len(current) >= max_b or cur_chars + blen > max_c) and current:
            groups.append(current)
            current   = []
            cur_chars = 0
        current.append(item)
        cur_chars += blen

    if current:
        groups.append(current)

    return groups if groups else [[]]


def split_code(lines, max_lines=MAX_CODE_LINES):
    """Split a list of code lines into groups that fit within one slide."""
    return [lines[i:i+max_lines] for i in range(0, len(lines), max_lines)] or [[]]


def add_content_slide(prs, title: str, bullet_items, contd_num: int = 0):
    """
    Add a standard bullet-content slide.
    bullet_items: list of {'type':'BODY'|'CODE', 'text': str}
    """
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Slide background
    set_slide_bg(slide, C_BODY_BG)

    # Title
    full_title = title
    if contd_num > 0:
        full_title = f'{title} (Contd. {contd_num})'
    add_title_box(slide, full_title)

    # Body text box
    body_box = slide.shapes.add_textbox(
        BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT
    )
    tf = body_box.text_frame
    tf.word_wrap = True

    first_para = True
    for item in bullet_items:
        text = item['text']
        if first_para:
            para = tf.paragraphs[0]
            first_para = False
        else:
            para = tf.add_paragraph()

        para.alignment = PP_ALIGN.LEFT
        para.space_after  = Pt(4)
        para.space_before = Pt(2)

        run = para.add_run()
        run.text = f'• {text}'
        run.font.size  = Pt(14)
        run.font.bold  = False
        run.font.color.rgb = C_BODY_TEXT
        run.font.name  = 'Calibri'


def add_code_slide(prs, title: str, code_lines, contd_num: int = 0):
    """
    Add a code-block slide with dark background and monospace font.
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Dark background for code slides
    set_slide_bg(slide, C_CODE_BG)

    # Title bar (lighter for code slides)
    shape = slide.shapes.add_shape(
        1,
        TITLE_LEFT, TITLE_TOP, TITLE_WIDTH, TITLE_HEIGHT
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_ACCENT
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    full_title = title
    if contd_num > 0:
        full_title = f'{title} (Contd. {contd_num})'
    run.text = full_title + '  [Code]'
    run.font.size  = Pt(20)
    run.font.bold  = True
    run.font.color.rgb = C_TITLE_TEXT
    run.font.name  = 'Calibri'

    # Code body
    code_box = slide.shapes.add_textbox(
        BODY_LEFT, BODY_TOP, BODY_WIDTH, BODY_HEIGHT
    )
    tf = code_box.text_frame
    tf.word_wrap = True

    first_para = True
    for line in code_lines:
        if first_para:
            para = tf.paragraphs[0]
            first_para = False
        else:
            para = tf.add_paragraph()
        run = para.add_run()
        run.text = line
        run.font.size  = Pt(13)
        run.font.bold  = False
        run.font.color.rgb = C_CODE_TEXT
        run.font.name  = 'Courier New'


def add_divider_slide(prs, title: str, level: int = 1):
    """Section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg_color = C_DIVIDER_BG if level == 1 else C_TITLE_BG
    set_slide_bg(slide, bg_color)

    fs = Pt(36) if level == 1 else Pt(28)
    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.8), Inches(11.33), Inches(2)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size  = fs
    run.font.bold  = True
    run.font.color.rgb = C_TITLE_TEXT
    run.font.name  = 'Calibri'


def add_cover_slide(prs, title: str):
    """Opening cover slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_bg(slide, C_TITLE_BG)

    txBox = slide.shapes.add_textbox(
        Inches(1), Inches(2.2), Inches(11.33), Inches(1.8)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size  = Pt(40)
    run.font.bold  = True
    run.font.color.rgb = C_TITLE_TEXT
    run.font.name  = 'Calibri'

    sub = tf.add_paragraph()
    sub.alignment = PP_ALIGN.CENTER
    r2 = sub.add_run()
    r2.text = 'Customer Information Control System'
    r2.font.size  = Pt(22)
    r2.font.bold  = False
    r2.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)
    r2.font.name  = 'Calibri'


# ── Section rendering ─────────────────────────────────────────────────────────
def render_section(prs, section: Section):
    """
    Render one section into the presentation.
    Handles bullet splitting and code splitting.
    """
    title = section.title

    # Add divider slide for the section
    add_divider_slide(prs, title, level=section.level)

    if not section.items:
        return

    # Group consecutive CODE lines together
    groups = []  # list of {'kind': 'bullets'|'code', 'items': [...]}
    current_kind  = None
    current_group = []

    for item in section.items:
        kind = 'code' if item['type'] == 'CODE' else 'bullets'
        if kind != current_kind and current_group:
            groups.append({'kind': current_kind, 'items': current_group})
            current_group = []
        current_kind = kind
        current_group.append(item)

    if current_group:
        groups.append({'kind': current_kind, 'items': current_group})

    slide_counter = 0  # counts content slides for this section (for Contd.)

    for grp in groups:
        if grp['kind'] == 'bullets':
            chunks = split_bullets(grp['items'])
            for chunk in chunks:
                slide_counter += 1
                add_content_slide(
                    prs, title, chunk,
                    contd_num=slide_counter - 1  # 0 = no suffix
                )
        else:
            # code group: extract lines
            lines = [it['text'] for it in grp['items']]
            chunks = split_code(lines)
            for chunk in chunks:
                slide_counter += 1
                add_code_slide(
                    prs, title, chunk,
                    contd_num=slide_counter - 1
                )


# ── Table handling ────────────────────────────────────────────────────────────
def add_table_slide(prs, title: str, rows, contd_num: int = 0):
    """Render a table as a text-based slide (simple representation)."""
    # Convert rows to bullet-like text, splitting long cells at sentences
    bullet_items = []
    for row in rows:
        cells = [c for c in row if c]
        if not cells:
            continue
        # Join columns with pipe separator
        line = ' | '.join(cells)
        # If a single cell is very long, split into multiple items
        if len(line) > MAX_CHARS:
            # Split by sentence or at a reasonable break
            sub_lines = textwrap.wrap(line, width=90)
            for sub in sub_lines:
                bullet_items.append({'type': 'BODY', 'text': sub})
        else:
            bullet_items.append({'type': 'BODY', 'text': line})

    if not bullet_items:
        return

    chunks = split_bullets(bullet_items)
    for i, chunk in enumerate(chunks):
        add_content_slide(prs, title, chunk, contd_num=i)


# ── Main conversion ───────────────────────────────────────────────────────────
def convert(input_path: str, output_path: str):
    print(f'Reading {input_path} …')
    doc = Document(input_path)

    print(f'  → {len(doc.paragraphs)} paragraphs, {len(doc.tables)} tables')

    # Classify paragraphs
    rows = classify_paragraphs(doc)

    # Build sections
    sections = build_sections(rows)
    print(f'  → {len(sections)} sections identified')

    # Extract tables (append after main content)
    tables = extract_tables(doc)

    # Create presentation (widescreen 16:9)
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Cover slide
    add_cover_slide(prs, 'CICS')

    # Render all sections
    for sec in sections:
        render_section(prs, sec)

    # Render tables at the end as supplementary slides
    if tables:
        add_divider_slide(prs, 'Reference Tables', level=1)
        for i, tbl in enumerate(tables, 1):
            add_table_slide(prs, f'Table {i}', tbl)

    prs.save(output_path)
    slide_count = len(prs.slides)
    print(f'Saved {output_path} ({slide_count} slides)')


if __name__ == '__main__':
    base = os.path.dirname(os.path.abspath(__file__))
    convert(
        os.path.join(base, 'CICS', 'CICS-new.docx'),
        os.path.join(base, 'CICS', 'CICS_Presentation.pptx'),
    )
